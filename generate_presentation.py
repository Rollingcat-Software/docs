"""
FIVUCSAS Presentation Generator
Generates a professional PowerPoint presentation for the January 7, 2026 defense.

Requirements:
    pip install python-pptx Pillow

Usage:
    python generate_presentation.py

Output:
    FIVUCSAS_Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import os

# ============================================================================
# CONFIGURATION
# ============================================================================

# Colors (FIVUCSAS Branding)
PRIMARY_COLOR = RgbColor(0x1E, 0x3A, 0x5F)      # Dark Blue #1e3a5f
ACCENT_COLOR = RgbColor(0xF3, 0x9C, 0x12)       # Orange #f39c12
WHITE = RgbColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RgbColor(0xF5, 0xF5, 0xF5)
DARK_GRAY = RgbColor(0x33, 0x33, 0x33)
SUCCESS_GREEN = RgbColor(0x27, 0xAE, 0x60)
CHECK_GREEN = RgbColor(0x2E, 0xCC, 0x71)

# Slide dimensions (16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Font settings
TITLE_FONT = "Segoe UI"
BODY_FONT = "Segoe UI"
TITLE_SIZE = Pt(36)
SUBTITLE_SIZE = Pt(24)
BODY_SIZE = Pt(18)
SMALL_SIZE = Pt(14)
PAGE_NUM_SIZE = Pt(12)

# Margins
LEFT_MARGIN = Inches(0.6)
TOP_MARGIN = Inches(0.8)
CONTENT_WIDTH = Inches(12.1)

TOTAL_SLIDES = 17


def create_presentation():
    """Create the FIVUCSAS presentation."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Use blank layout
    blank_layout = prs.slide_layouts[6]

    # Create all slides
    create_slide_1_title(prs, blank_layout)
    create_slide_2_outline(prs, blank_layout)
    create_slide_3_problem(prs, blank_layout)
    create_slide_4_related_work(prs, blank_layout)
    create_slide_5_scope(prs, blank_layout)
    create_slide_6_architecture(prs, blank_layout)
    create_slide_7_biometric_puzzle(prs, blank_layout)
    create_slide_8_biometric_demo(prs, blank_layout)
    create_slide_9_ml_pipeline(prs, blank_layout)
    create_slide_10_card_nfc(prs, blank_layout)
    create_slide_11_doc_demo(prs, blank_layout)
    create_slide_12_tasks(prs, blank_layout)
    create_slide_13_challenges(prs, blank_layout)
    create_slide_14_future(prs, blank_layout)
    create_slide_15_thank_you(prs, blank_layout)
    create_slide_16_references(prs, blank_layout)
    create_slide_17_qa(prs, blank_layout)

    return prs


# ============================================================================
# HELPER FUNCTIONS
# ============================================================================

def add_background(slide, color=PRIMARY_COLOR):
    """Add solid background color to slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_white_background(slide):
    """Add white background to slide."""
    add_background(slide, WHITE)


def add_title(slide, text, top=TOP_MARGIN, color=PRIMARY_COLOR, size=TITLE_SIZE):
    """Add title text to slide."""
    title_box = slide.shapes.add_textbox(LEFT_MARGIN, top, CONTENT_WIDTH, Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = TITLE_FONT
    p.font.size = size
    p.font.bold = True
    p.font.color.rgb = color
    return title_box


def add_subtitle_line(slide, top):
    """Add an accent line under the title."""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        LEFT_MARGIN, top, Inches(2), Pt(4)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_COLOR
    line.line.fill.background()
    return line


def add_page_number(slide, page_num):
    """Add page number to bottom right."""
    page_box = slide.shapes.add_textbox(
        Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.3)
    )
    tf = page_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{page_num}/{TOTAL_SLIDES}"
    p.font.name = BODY_FONT
    p.font.size = PAGE_NUM_SIZE
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.RIGHT


def add_text_box(slide, text, left, top, width, height, font_size=BODY_SIZE,
                 color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT):
    """Add a text box with specified properties."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = BODY_FONT
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return box


def add_bullet_list(slide, items, left, top, width, height, font_size=BODY_SIZE):
    """Add a bulleted list."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.name = BODY_FONT
        p.font.size = font_size
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(8)
    return box


def add_stat_box(slide, value, label, left, top, width=Inches(2.5), height=Inches(1.8)):
    """Add a statistics box with large number and label."""
    # Box background
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_GRAY
    box.line.fill.background()

    # Value text
    value_box = slide.shapes.add_textbox(left, top + Inches(0.2), width, Inches(0.8))
    tf = value_box.text_frame
    p = tf.paragraphs[0]
    p.text = value
    p.font.name = TITLE_FONT
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    p.alignment = PP_ALIGN.CENTER

    # Label text
    label_box = slide.shapes.add_textbox(left, top + Inches(1.0), width, Inches(0.6))
    tf = label_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.name = BODY_FONT
    p.font.size = SMALL_SIZE
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER


def add_content_box(slide, title, content_lines, left, top, width, height):
    """Add a content box with title and bullet points."""
    # Box background
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_GRAY
    box.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.4))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = BODY_FONT
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR

    # Content
    content_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.5), width - Inches(0.4), height - Inches(0.6))
    tf = content_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {line}"
        p.font.name = BODY_FONT
        p.font.size = SMALL_SIZE
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(4)


# ============================================================================
# SLIDE CREATION FUNCTIONS
# ============================================================================

def create_slide_1_title(prs, layout):
    """SLIDE 1: Title Slide"""
    slide = prs.slides.add_slide(layout)
    add_background(slide, PRIMARY_COLOR)

    # Main title
    title_box = slide.shapes.add_textbox(LEFT_MARGIN, Inches(1.8), CONTENT_WIDTH, Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "FIVUCSAS"
    p.font.name = TITLE_FONT
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(LEFT_MARGIN, Inches(3.0), CONTENT_WIDTH, Inches(0.8))
    tf = subtitle_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Face and Identity Verification Using Cloud-Based SaaS Models"
    p.font.name = BODY_FONT
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Course info
    course_box = slide.shapes.add_textbox(LEFT_MARGIN, Inches(3.8), CONTENT_WIDTH, Inches(0.5))
    tf = course_box.text_frame
    p = tf.paragraphs[0]
    p.text = "CSE4197 Engineering Project - Fall 2025"
    p.font.name = BODY_FONT
    p.font.size = Pt(18)
    p.font.color.rgb = ACCENT_COLOR
    p.alignment = PP_ALIGN.CENTER

    # Team members
    team_box = slide.shapes.add_textbox(LEFT_MARGIN, Inches(4.6), CONTENT_WIDTH, Inches(1.2))
    tf = team_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Team:"
    p.font.name = BODY_FONT
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    members = ["Ahmet Abdullah Gultekin (150121025)",
               "Ayse Gulsum Eren (150120005)",
               "Aysenur Arici (150123825)"]
    for member in members:
        p = tf.add_paragraph()
        p.text = member
        p.font.name = BODY_FONT
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # Advisor
    advisor_box = slide.shapes.add_textbox(LEFT_MARGIN, Inches(6.0), CONTENT_WIDTH, Inches(0.5))
    tf = advisor_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Supervisor: Assoc. Prof. Dr. Mustafa Agaoglu"
    p.font.name = BODY_FONT
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # University
    uni_box = slide.shapes.add_textbox(LEFT_MARGIN, Inches(6.5), CONTENT_WIDTH, Inches(0.4))
    tf = uni_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Marmara University - Computer Engineering"
    p.font.name = BODY_FONT
    p.font.size = Pt(14)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER

    # Page number (white for dark background)
    page_box = slide.shapes.add_textbox(Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.3))
    tf = page_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"1/{TOTAL_SLIDES}"
    p.font.name = BODY_FONT
    p.font.size = PAGE_NUM_SIZE
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.RIGHT


def create_slide_2_outline(prs, layout):
    """SLIDE 2: Presentation Outline"""
    slide = prs.slides.add_slide(layout)
    add_white_background(slide)

    add_title(slide, "OUTLINE")
    add_subtitle_line(slide, Inches(1.4))

    sections = [
        ("WHY THIS MATTERS", [
            "1. Problem Statement & Motivation",
            "2. Related Work & Gap Analysis",
            "3. Scope & Engineering Constraints"
        ]),
        ("HOW WE SOLVE IT", [
            "4. System Architecture",
            "5. The Biometric Puzzle (Hybrid Liveness)",
            "6. ML Pipeline & Vector Search",
            "7. Card Detection & NFC Verification"
        ]),
        ("WHAT WE BUILT", [
            "8. Tasks Accomplished",
            "9. Technical Challenges & Solutions",
            "10. Future Work & Contingency Plan"
        ])
    ]

    top = Inches(1.8)
    for section_title, items in sections:
        # Section header
        header_box = slide.shapes.add_textbox(LEFT_MARGIN, top, CONTENT_WIDTH, Inches(0.4))
        tf = header_box.text_frame
        p = tf.paragraphs[0]
        p.text = section_title
        p.font.name = BODY_FONT
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR

        top += Inches(0.5)

        # Items
        for item in items:
            item_box = slide.shapes.add_textbox(LEFT_MARGIN + Inches(0.3), top, CONTENT_WIDTH, Inches(0.35))
            tf = item_box.text_frame
            p = tf.paragraphs[0]
            p.text = item
            p.font.name = BODY_FONT
            p.font.size = BODY_SIZE
            p.font.color.rgb = DARK_GRAY
            top += Inches(0.4)

        top += Inches(0.3)

    add_page_number(slide, 2)


def create_slide_3_problem(prs, layout):
    """SLIDE 3: Problem Statement & Motivation"""
    slide = prs.slides.add_slide(layout)
    add_white_background(slide)

    add_title(slide, "WHY THIS MATTERS")
    add_subtitle_line(slide, Inches(1.4))

    # Quote box
    quote_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        LEFT_MARGIN, Inches(1.7), CONTENT_WIDTH, Inches(0.8)
    )
    quote_box.fill.solid()
    quote_box.fill.fore_color.rgb = RgbColor(0xFF, 0xF3, 0xCD)  # Light yellow
    quote_box.line.fill.background()

    quote_text = slide.shapes.add_textbox(LEFT_MARGIN + Inches(0.3), Inches(1.9), CONTENT_WIDTH - Inches(0.6), Inches(0.5))
    tf = quote_text.text_frame
    p = tf.paragraphs[0]
    p.text = '"2024: Deepfake CFO video call -> $25 Million stolen"'
    p.font.name = BODY_FONT
    p.font.size = Pt(20)
    p.font.italic = True
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER

    # Statistics boxes
    add_stat_box(slide, "$23B", "Identity Fraud\nLosses (2024)", Inches(1.5), Inches(3.0))
    add_stat_box(slide, "+400%", "Deepfake\nAttacks (YoY)", Inches(5.4), Inches(3.0))
    add_stat_box(slide, "1 in 4", "People Cannot\nDetect Deepfakes", Inches(9.3), Inches(3.0))

    # Goal statement
    goal_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.5), Inches(5.5), Inches(10.3), Inches(1.0)
    )
    goal_box.fill.solid()
    goal_box.fill.fore_color.rgb = PRIMARY_COLOR
    goal_box.line.fill.background()

    goal_text = slide.shapes.add_textbox(Inches(1.5), Inches(5.7), Inches(10.3), Inches(0.6))
    tf = goal_text.text_frame
    p = tf.paragraphs[0]
    p.text = "OUR GOAL: Prove LIVE + AUTHENTIC + IMPOSSIBLE to spoof"
    p.font.name = BODY_FONT
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    add_page_number(slide, 3)


def create_slide_4_related_work(prs, layout):
    """SLIDE 4: Related Work & Gap Analysis"""
    slide = prs.slides.add_slide(layout)
    add_white_background(slide)

    add_title(slide, "RELATED WORK & GAP ANALYSIS")
    add_subtitle_line(slide, Inches(1.4))

    # Table header
    headers = ["Feature", "Azure", "AWS", "Sodec", "BioGATE", "FIVUCSAS"]
    col_widths = [Inches(2.5), Inches(1.3), Inches(1.3), Inches(1.3), Inches(1.3), Inches(1.8)]

    table_left = Inches(1.0)
    table_top = Inches(1.8)
    row_height = Inches(0.5)

    # Header row
    x = table_left
    for i, header in enumerate(headers):
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, table_top, col_widths[i], row_height)
        box.fill.solid()
        box.fill.fore_color.rgb = PRIMARY_COLOR
        box.line.color.rgb = WHITE

        text_box = slide.shapes.add_textbox(x, table_top + Inches(0.1), col_widths[i], row_height)
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = header
        p.font.name = BODY_FONT
        p.font.size = SMALL_SIZE
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        x += col_widths[i]

    # Data rows
    data = [
        ["Open Source", "X", "X", "X", "X", "CHECK"],
        ["Liveness", "CHECK", "CHECK", "X", "CHECK", "CHECK"],
        ["Multi-Tenant", "CHECK", "CHECK", "X", "X", "CHECK"],
        ["Multi-Platform", "CHECK", "CHECK", "CHECK", "X", "CHECK"],
        ["NFC ICAO", "X", "X", "X", "CHECK", "CHECK"],
        ["Card Detection", "X", "X", "X", "X", "CHECK"],
        ["Hybrid Liveness", "X", "X", "X", "X", "CHECK"],
    ]

    for row_idx, row in enumerate(data):
        x = table_left
        y = table_top + row_height * (row_idx + 1)

        for col_idx, cell in enumerate(row):
            bg_color = LIGHT_GRAY if row_idx % 2 == 0 else WHITE
            if col_idx == 5:  # FIVUCSAS column
                bg_color = RgbColor(0xE8, 0xF5, 0xE9)  # Light green

            box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, col_widths[col_idx], row_height)
            box.fill.solid()
            box.fill.fore_color.rgb = bg_color
            box.line.color.rgb = LIGHT_GRAY

            text_box = slide.shapes.add_textbox(x, y + Inches(0.1), col_widths[col_idx], row_height)
            tf = text_box.text_frame
            p = tf.paragraphs[0]

            if cell == "CHECK":
                p.text = "CHECK"  # Checkmark
                p.font.color.rgb = SUCCESS_GREEN
            elif cell == "X":
                p.text = "X"  # X mark
                p.font.color.rgb = RgbColor(0xE7, 0x4C, 0x3C)  # Red
            else:
                p.text = cell
                p.font.color.rgb = DARK_GRAY

            p.font.name = BODY_FONT
            p.font.size = SMALL_SIZE
            p.font.bold = (col_idx == 0)
            p.alignment = PP_ALIGN.CENTER
            x += col_widths[col_idx]

    # Gap statement
    gap_box = slide.shapes.add_textbox(LEFT_MARGIN, Inches(6.0), CONTENT_WIDTH, Inches(0.5))
    tf = gap_box.text_frame
    p = tf.paragraphs[0]
    p.text = "GAP: No open-source solution with hybrid liveness + NFC + card detection"
    p.font.name = BODY_FONT
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    p.alignment = PP_ALIGN.CENTER

    add_page_number(slide, 4)


def create_slide_5_scope(prs, layout):
    """SLIDE 5: Scope & Engineering Constraints"""
    slide = prs.slides.add_slide(layout)
    add_white_background(slide)

    add_title(slide, "SCOPE & ENGINEERING CONSTRAINTS")
    add_subtitle_line(slide, Inches(1.4))

    # In Scope section
    in_scope_header = slide.shapes.add_textbox(LEFT_MARGIN, Inches(1.8), Inches(5.5), Inches(0.4))
    tf = in_scope_header.text_frame
    p = tf.paragraphs[0]
    p.text = "IN SCOPE (MVP)"
    p.font.name = BODY_FONT
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = SUCCESS_GREEN

    in_scope_items = [
        "Cloud-Native SaaS Platform",
        "Hybrid Liveness Detection",
        "Card Detection (ML Model)",
        "NFC Document Reading",
        "Multi-Tenant Admin Dashboard",
        "Cross-Platform Client Apps"
    ]

    for i, item in enumerate(in_scope_items):
        item_box = slide.shapes.add_textbox(LEFT_MARGIN + Inches(0.2), Inches(2.3 + i * 0.4), Inches(5.5), Inches(0.35))
        tf = item_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"CHECK {item}"
        p.font.name = BODY_FONT
        p.font.size = BODY_SIZE
        p.font.color.rgb = DARK_GRAY

    # Out of Scope section
    out_scope_header = slide.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.5), Inches(0.4))
    tf = out_scope_header.text_frame
    p = tf.paragraphs[0]
    p.text = "OUT OF SCOPE"
    p.font.name = BODY_FONT
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RgbColor(0xE7, 0x4C, 0x3C)  # Red

    out_scope_item = slide.shapes.add_textbox(Inches(7.2), Inches(2.3), Inches(5.5), Inches(0.35))
    tf = out_scope_item.text_frame
    p = tf.paragraphs[0]
    p.text = "X Hardware Manufacturing"
    p.font.name = BODY_FONT
    p.font.size = BODY_SIZE
    p.font.color.rgb = DARK_GRAY

    # Engineering Constraints box
    constraints_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        LEFT_MARGIN, Inches(5.0), CONTENT_WIDTH, Inches(1.5)
    )
    constraints_box.fill.solid()
    constraints_box.fill.fore_color.rgb = LIGHT_GRAY
    constraints_box.line.fill.background()

    constraints_title = slide.shapes.add_textbox(LEFT_MARGIN + Inches(0.3), Inches(5.1), CONTENT_WIDTH, Inches(0.4))
    tf = constraints_title.text_frame
    p = tf.paragraphs[0]
    p.text = "ENGINEERING CONSTRAINTS"
    p.font.name = BODY_FONT
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR

    constraints = [
        ("> 480p", "Image Quality"),
        ("< 200ms", "API Latency"),
        ("ISO 14443", "NFC Standard")
    ]

    x_positions = [Inches(1.5), Inches(5.5), Inches(9.5)]
    for i, (value, label) in enumerate(constraints):
        value_box = slide.shapes.add_textbox(x_positions[i], Inches(5.6), Inches(3.0), Inches(0.5))
        tf = value_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.name = TITLE_FONT
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR
        p.alignment = PP_ALIGN.CENTER

        label_box = slide.shapes.add_textbox(x_positions[i], Inches(6.1), Inches(3.0), Inches(0.3))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.name = BODY_FONT
        p.font.size = SMALL_SIZE
        p.font.color.rgb = DARK_GRAY
        p.alignment = PP_ALIGN.CENTER

    add_page_number(slide, 5)


def create_slide_6_architecture(prs, layout):
    """SLIDE 6: System Architecture"""
    slide = prs.slides.add_slide(layout)
    add_white_background(slide)

    add_title(slide, "SYSTEM ARCHITECTURE")
    add_subtitle_line(slide, Inches(1.4))

    # API Gateway box
    gateway_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(4.5), Inches(1.7), Inches(4.3), Inches(0.9)
    )
    gateway_box.fill.solid()
    gateway_box.fill.fore_color.rgb = ACCENT_COLOR
    gateway_box.line.fill.background()

    gateway_text = slide.shapes.add_textbox(Inches(4.5), Inches(1.8), Inches(4.3), Inches(0.7))
    tf = gateway_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "NGINX API GATEWAY"
    p.font.name = BODY_FONT
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = "Rate Limiting | Routing | Load Balancing"
    p.font.name = BODY_FONT
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Three service boxes
    services = [
        ("IDENTITY CORE", "Spring Boot (Java 21)", ["JWT Auth", "Multi-Tenant", "RBAC"]),
        ("BIOMETRIC PROC.", "FastAPI (Python 3.11)", ["Face Detection", "40+ Endpoints", "Liveness"]),
        ("CLIENT APPS", "Kotlin Multiplatform", ["Android/iOS/Desktop", "Camera/NFC", "Card Detect"])
    ]

    x_positions = [Inches(0.8), Inches(4.8), Inches(8.8)]
    for i, (title, subtitle, features) in enumerate(services):
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x_positions[i], Inches(3.0), Inches(3.7), Inches(2.2)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GRAY
        box.line.fill.background()

        title_box = slide.shapes.add_textbox(x_positions[i], Inches(3.1), Inches(3.7), Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = BODY_FONT
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR
        p.alignment = PP_ALIGN.CENTER

        subtitle_box = slide.shapes.add_textbox(x_positions[i], Inches(3.5), Inches(3.7), Inches(0.3))
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.name = BODY_FONT
        p.font.size = Pt(11)
        p.font.color.rgb = ACCENT_COLOR
        p.alignment = PP_ALIGN.CENTER

        for j, feature in enumerate(features):
            feature_box = slide.shapes.add_textbox(x_positions[i] + Inches(0.2), Inches(4.0 + j * 0.35), Inches(3.3), Inches(0.3))
            tf = feature_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"• {feature}"
            p.font.name = BODY_FONT
            p.font.size = Pt(12)
            p.font.color.rgb = DARK_GRAY

    # Database box
    db_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.5), Inches(5.6), Inches(10.3), Inches(1.0)
    )
    db_box.fill.solid()
    db_box.fill.fore_color.rgb = PRIMARY_COLOR
    db_box.line.fill.background()

    db_text = slide.shapes.add_textbox(Inches(1.5), Inches(5.75), Inches(5.0), Inches(0.6))
    tf = db_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "PostgreSQL 16 + pgvector"
    p.font.name = BODY_FONT
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p = tf.add_paragraph()
    p.text = "Vector Embeddings | IVFFlat Indexing"
    p.font.name = BODY_FONT
    p.font.size = Pt(11)
    p.font.color.rgb = LIGHT_GRAY

    redis_text = slide.shapes.add_textbox(Inches(7.0), Inches(5.75), Inches(4.5), Inches(0.6))
    tf = redis_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Redis"
    p.font.name = BODY_FONT
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p = tf.add_paragraph()
    p.text = "Cache | Event Bus"
    p.font.name = BODY_FONT
    p.font.size = Pt(11)
    p.font.color.rgb = LIGHT_GRAY

    # Architecture pattern note
    pattern_box = slide.shapes.add_textbox(LEFT_MARGIN, Inches(6.8), CONTENT_WIDTH, Inches(0.3))
    tf = pattern_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Architecture: Hexagonal (Ports & Adapters) + DDD"
    p.font.name = BODY_FONT
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR

    add_page_number(slide, 6)


def create_slide_7_biometric_puzzle(prs, layout):
    """SLIDE 7: The Biometric Puzzle (Liveness Detection)"""
    slide = prs.slides.add_slide(layout)
    add_white_background(slide)

    add_title(slide, "THE BIOMETRIC PUZZLE - HYBRID LIVENESS")
    add_subtitle_line(slide, Inches(1.4))

    # Flow diagram - Server side
    server_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(1.9), Inches(3.5), Inches(3.5)
    )
    server_box.fill.solid()
    server_box.fill.fore_color.rgb = PRIMARY_COLOR
    server_box.line.fill.background()

    server_title = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(3.5), Inches(0.4))
    tf = server_title.text_frame
    p = tf.paragraphs[0]
    p.text = "SERVER"
    p.font.name = BODY_FONT
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    server_steps = [
        "1. Generate Challenge",
        '   "Blink Left Eye"',
        "",
        "2. Calculate EAR",
        "   EAR = 0.18 < 0.2",
        "   CHECKMARK Blink Detected",
        "",
        "3. Passive Analysis",
        "   LBP Texture Check",
        "   CHECKMARK Not a Screen"
    ]

    server_content = slide.shapes.add_textbox(Inches(0.9), Inches(2.5), Inches(3.3), Inches(2.8))
    tf = server_content.text_frame
    tf.word_wrap = True
    for i, step in enumerate(server_steps):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = step
        p.font.name = BODY_FONT
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE if "CHECKMARK" not in step else CHECK_GREEN

    # Client side
    client_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(9.0), Inches(1.9), Inches(3.5), Inches(3.5)
    )
    client_box.fill.solid()
    client_box.fill.fore_color.rgb = LIGHT_GRAY
    client_box.line.fill.background()

    client_title = slide.shapes.add_textbox(Inches(9.0), Inches(2.0), Inches(3.5), Inches(0.4))
    tf = client_title.text_frame
    p = tf.paragraphs[0]
    p.text = "MOBILE CLIENT"
    p.font.name = BODY_FONT
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    p.alignment = PP_ALIGN.CENTER

    client_steps = [
        "Camera Capture",
        "",
        "468 Landmarks",
        "(MediaPipe)",
        "",
        "Real-time",
        "Processing",
        "",
        "CHECKMARK PASS"
    ]

    client_content = slide.shapes.add_textbox(Inches(9.1), Inches(2.5), Inches(3.3), Inches(2.8))
    tf = client_content.text_frame
    tf.word_wrap = True
    for i, step in enumerate(client_steps):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = step
        p.font.name = BODY_FONT
        p.font.size = Pt(13)
        p.font.color.rgb = SUCCESS_GREEN if "CHECKMARK" in step else DARK_GRAY
        p.alignment = PP_ALIGN.CENTER

    # Arrows between boxes (using text for simplicity)
    arrow1 = slide.shapes.add_textbox(Inches(4.5), Inches(2.5), Inches(4.3), Inches(0.4))
    tf = arrow1.text_frame
    p = tf.paragraphs[0]
    p.text = "Challenge >>>>"
    p.font.name = BODY_FONT
    p.font.size = Pt(14)
    p.font.color.rgb = ACCENT_COLOR
    p.alignment = PP_ALIGN.CENTER

    arrow2 = slide.shapes.add_textbox(Inches(4.5), Inches(3.8), Inches(4.3), Inches(0.4))
    tf = arrow2.text_frame
    p = tf.paragraphs[0]
    p.text = "<<<< Response"
    p.font.name = BODY_FONT
    p.font.size = Pt(14)
    p.font.color.rgb = ACCENT_COLOR
    p.alignment = PP_ALIGN.CENTER

    # Detection methods
    active_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        LEFT_MARGIN, Inches(5.7), Inches(5.8), Inches(1.0)
    )
    active_box.fill.solid()
    active_box.fill.fore_color.rgb = RgbColor(0xE3, 0xF2, 0xFD)  # Light blue
    active_box.line.fill.background()

    active_title = slide.shapes.add_textbox(LEFT_MARGIN + Inches(0.2), Inches(5.8), Inches(5.4), Inches(0.3))
    tf = active_title.text_frame
    p = tf.paragraphs[0]
    p.text = "ACTIVE DETECTION"
    p.font.name = BODY_FONT
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR

    active_content = slide.shapes.add_textbox(LEFT_MARGIN + Inches(0.2), Inches(6.1), Inches(5.4), Inches(0.4))
    tf = active_content.text_frame
    p = tf.paragraphs[0]
    p.text = "EAR (Eye Aspect Ratio), MAR (Mouth), Head Pose"
    p.font.name = BODY_FONT
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_GRAY

    passive_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.8), Inches(5.7), Inches(5.8), Inches(1.0)
    )
    passive_box.fill.solid()
    passive_box.fill.fore_color.rgb = RgbColor(0xFC, 0xE4, 0xEC)  # Light pink
    passive_box.line.fill.background()

    passive_title = slide.shapes.add_textbox(Inches(7.0), Inches(5.8), Inches(5.4), Inches(0.3))
    tf = passive_title.text_frame
    p = tf.paragraphs[0]
    p.text = "PASSIVE DETECTION"
    p.font.name = BODY_FONT
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR

    passive_content = slide.shapes.add_textbox(Inches(7.0), Inches(6.1), Inches(5.4), Inches(0.4))
    tf = passive_content.text_frame
    p = tf.paragraphs[0]
    p.text = "LBP Texture, Color Distribution, Frequency Domain"
    p.font.name = BODY_FONT
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_GRAY

    add_page_number(slide, 7)


def create_slide_8_biometric_demo(prs, layout):
    """SLIDE 8: Biometric Processor Demo"""
    slide = prs.slides.add_slide(layout)
    add_white_background(slide)

    add_title(slide, "BIOMETRIC PROCESSOR IN ACTION")
    add_subtitle_line(slide, Inches(1.4))

    # Three demo boxes for team members
    team_members = [
        ("AHMET", "EAR: 0.28", "Quality: 94%"),
        ("AYSENUR", "EAR: 0.31", "Quality: 96%"),
        ("GULSUM", "EAR: 0.29", "Quality: 95%")
    ]

    x_positions = [Inches(1.0), Inches(4.8), Inches(8.6)]

    for i, (name, ear, quality) in enumerate(team_members):
        # Main photo placeholder box
        photo_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x_positions[i], Inches(1.9), Inches(3.5), Inches(3.2)
        )
        photo_box.fill.solid()
        photo_box.fill.fore_color.rgb = LIGHT_GRAY
        photo_box.line.color.rgb = PRIMARY_COLOR
        photo_box.line.width = Pt(2)

        # Placeholder text
        placeholder = slide.shapes.add_textbox(x_positions[i], Inches(2.8), Inches(3.5), Inches(1.0))
        tf = placeholder.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"[{name}'S FACE"
        p.font.name = BODY_FONT
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.alignment = PP_ALIGN.CENTER
        p = tf.add_paragraph()
        p.text = "with 468 MESH"
        p.font.name = BODY_FONT
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.alignment = PP_ALIGN.CENTER
        p = tf.add_paragraph()
        p.text = "OVERLAY]"
        p.font.name = BODY_FONT
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.alignment = PP_ALIGN.CENTER

        # Metrics box
        metrics_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x_positions[i] + Inches(0.3), Inches(4.3), Inches(2.9), Inches(0.9)
        )
        metrics_box.fill.solid()
        metrics_box.fill.fore_color.rgb = WHITE
        metrics_box.line.color.rgb = SUCCESS_GREEN
        metrics_box.line.width = Pt(2)

        metrics_text = slide.shapes.add_textbox(x_positions[i] + Inches(0.4), Inches(4.4), Inches(2.7), Inches(0.7))
        tf = metrics_text.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = ear
        p.font.name = BODY_FONT
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_GRAY
        p = tf.add_paragraph()
        p.text = quality
        p.font.name = BODY_FONT
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_GRAY
        p = tf.add_paragraph()
        p.text = "CHECKMARK LIVE"
        p.font.name = BODY_FONT
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = SUCCESS_GREEN

    # Caption
    caption = slide.shapes.add_textbox(LEFT_MARGIN, Inches(5.8), CONTENT_WIDTH, Inches(0.6))
    tf = caption.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "All team members verified with 468 facial landmarks + liveness detection"
    p.font.name = BODY_FONT
    p.font.size = Pt(16)
    p.font.color.rgb = PRIMARY_COLOR
    p.alignment = PP_ALIGN.CENTER

    # Note box
    note_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2.0), Inches(6.3), Inches(9.3), Inches(0.6)
    )
    note_box.fill.solid()
    note_box.fill.fore_color.rgb = RgbColor(0xFF, 0xF3, 0xCD)  # Light yellow
    note_box.line.fill.background()

    note_text = slide.shapes.add_textbox(Inches(2.2), Inches(6.4), Inches(8.9), Inches(0.4))
    tf = note_text.text_frame
    p = tf.paragraphs[0]
    p.text = "Note: Replace placeholder boxes with actual screenshots from web demo"
    p.font.name = BODY_FONT
    p.font.size = Pt(12)
    p.font.italic = True
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER

    add_page_number(slide, 8)


def create_slide_9_ml_pipeline(prs, layout):
    """SLIDE 9: ML Pipeline & Vector Search"""
    slide = prs.slides.add_slide(layout)
    add_white_background(slide)

    add_title(slide, "ML PIPELINE & VECTOR SEARCH")
    add_subtitle_line(slide, Inches(1.4))

    # Pipeline box
    pipeline_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        LEFT_MARGIN, Inches(1.8), CONTENT_WIDTH, Inches(2.5)
    )
    pipeline_box.fill.solid()
    pipeline_box.fill.fore_color.rgb = LIGHT_GRAY
    pipeline_box.line.fill.background()

    pipeline_title = slide.shapes.add_textbox(LEFT_MARGIN + Inches(0.2), Inches(1.9), CONTENT_WIDTH, Inches(0.4))
    tf = pipeline_title.text_frame
    p = tf.paragraphs[0]
    p.text = "RECOGNITION PIPELINE"
    p.font.name = BODY_FONT
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    p.alignment = PP_ALIGN.CENTER

    # Pipeline stages
    stages = [
        ("Input", "Image"),
        ("Detect", "MediaPipe"),
        ("Align", "5-Point"),
        ("Extract", "Embedding"),
        ("Search", "Cosine")
    ]

    stage_width = Inches(2.0)
    x_start = Inches(1.2)

    for i, (top_label, bottom_label) in enumerate(stages):
        x = x_start + i * Inches(2.4)

        # Stage box
        stage_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, Inches(2.6), stage_width, Inches(1.2)
        )
        stage_box.fill.solid()
        stage_box.fill.fore_color.rgb = PRIMARY_COLOR
        stage_box.line.fill.background()

        top_text = slide.shapes.add_textbox(x, Inches(2.7), stage_width, Inches(0.4))
        tf = top_text.text_frame
        p = tf.paragraphs[0]
        p.text = top_label
        p.font.name = BODY_FONT
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        bottom_text = slide.shapes.add_textbox(x, Inches(3.2), stage_width, Inches(0.4))
        tf = bottom_text.text_frame
        p = tf.paragraphs[0]
        p.text = bottom_label
        p.font.name = BODY_FONT
        p.font.size = Pt(12)
        p.font.color.rgb = LIGHT_GRAY
        p.alignment = PP_ALIGN.CENTER

        # Arrow (except for last)
        if i < len(stages) - 1:
            arrow = slide.shapes.add_textbox(x + stage_width, Inches(3.0), Inches(0.4), Inches(0.4))
            tf = arrow.text_frame
            p = tf.paragraphs[0]
            p.text = ">"
            p.font.name = BODY_FONT
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = ACCENT_COLOR

    # Vector Database box
    db_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        LEFT_MARGIN, Inches(4.6), CONTENT_WIDTH, Inches(1.2)
    )
    db_box.fill.solid()
    db_box.fill.fore_color.rgb = PRIMARY_COLOR
    db_box.line.fill.background()

    db_title = slide.shapes.add_textbox(LEFT_MARGIN + Inches(0.2), Inches(4.7), CONTENT_WIDTH, Inches(0.4))
    tf = db_title.text_frame
    p = tf.paragraphs[0]
    p.text = "VECTOR DATABASE"
    p.font.name = BODY_FONT
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    db_content = slide.shapes.add_textbox(LEFT_MARGIN + Inches(0.2), Inches(5.1), CONTENT_WIDTH, Inches(0.5))
    tf = db_content.text_frame
    p = tf.paragraphs[0]
    p.text = "PostgreSQL + pgvector + IVFFlat -> O(log n) sub-ms queries"
    p.font.name = BODY_FONT
    p.font.size = Pt(14)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER

    # Threshold info
    threshold_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(3.5), Inches(6.0), Inches(6.3), Inches(0.8)
    )
    threshold_box.fill.solid()
    threshold_box.fill.fore_color.rgb = SUCCESS_GREEN
    threshold_box.line.fill.background()

    threshold_text = slide.shapes.add_textbox(Inches(3.5), Inches(6.15), Inches(6.3), Inches(0.5))
    tf = threshold_text.text_frame
    p = tf.paragraphs[0]
    p.text = "Threshold: cosine distance < 0.68 = MATCH"
    p.font.name = BODY_FONT
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    add_page_number(slide, 9)


def create_slide_10_card_nfc(prs, layout):
    """SLIDE 10: Card Detection & NFC Verification"""
    slide = prs.slides.add_slide(layout)
    add_white_background(slide)

    add_title(slide, "CARD DETECTION & NFC VERIFICATION")
    add_subtitle_line(slide, Inches(1.4))

    # Step 1: Visual Card Detection
    step1_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        LEFT_MARGIN, Inches(1.8), CONTENT_WIDTH, Inches(2.0)
    )
    step1_box.fill.solid()
    step1_box.fill.fore_color.rgb = LIGHT_GRAY
    step1_box.line.fill.background()

    step1_title = slide.shapes.add_textbox(LEFT_MARGIN + Inches(0.2), Inches(1.9), CONTENT_WIDTH, Inches(0.4))
    tf = step1_title.text_frame
    p = tf.paragraphs[0]
    p.text = "STEP 1: VISUAL CARD DETECTION (ML MODEL)"
    p.font.name = BODY_FONT
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR

    step1_flow = slide.shapes.add_textbox(LEFT_MARGIN + Inches(0.2), Inches(2.4), CONTENT_WIDTH, Inches(0.4))
    tf = step1_flow.text_frame
    p = tf.paragraphs[0]
    p.text = "[Camera] -> [Trained Model] -> [Card Type Detected]"
    p.font.name = BODY_FONT
    p.font.size = Pt(14)
    p.font.color.rgb = ACCENT_COLOR
    p.alignment = PP_ALIGN.CENTER

    step1_features = [
        "Automatic ID type recognition (Turkish eID, Passport)",
        "Real-time detection via on-device ML",
        "Guides user for optimal card positioning"
    ]

    for i, feature in enumerate(step1_features):
        feature_box = slide.shapes.add_textbox(LEFT_MARGIN + Inches(0.3), Inches(2.9 + i * 0.35), CONTENT_WIDTH, Inches(0.3))
        tf = feature_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"• {feature}"
        p.font.name = BODY_FONT
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_GRAY

    # Step 2: NFC Chip Verification
    step2_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        LEFT_MARGIN, Inches(4.0), CONTENT_WIDTH, Inches(2.0)
    )
    step2_box.fill.solid()
    step2_box.fill.fore_color.rgb = LIGHT_GRAY
    step2_box.line.fill.background()

    step2_title = slide.shapes.add_textbox(LEFT_MARGIN + Inches(0.2), Inches(4.1), CONTENT_WIDTH, Inches(0.4))
    tf = step2_title.text_frame
    p = tf.paragraphs[0]
    p.text = "STEP 2: NFC CHIP VERIFICATION (ICAO)"
    p.font.name = BODY_FONT
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR

    step2_flow = slide.shapes.add_textbox(LEFT_MARGIN + Inches(0.2), Inches(4.6), CONTENT_WIDTH, Inches(0.4))
    tf = step2_flow.text_frame
    p = tf.paragraphs[0]
    p.text = "[Mobile NFC] -> [BAC Handshake] -> [Read DG1/DG2] -> [SOD CHECKMARK]"
    p.font.name = BODY_FONT
    p.font.size = Pt(14)
    p.font.color.rgb = ACCENT_COLOR
    p.alignment = PP_ALIGN.CENTER

    step2_features = [
        "MRZ-derived session keys (3DES secure messaging)",
        "DG1: Personal data, DG2: High-res JPEG2000 photo",
        "Digital signature proves document authenticity"
    ]

    for i, feature in enumerate(step2_features):
        feature_box = slide.shapes.add_textbox(LEFT_MARGIN + Inches(0.3), Inches(5.1 + i * 0.35), CONTENT_WIDTH, Inches(0.3))
        tf = feature_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"• {feature}"
        p.font.name = BODY_FONT
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_GRAY

    # Standards footer
    standards_box = slide.shapes.add_textbox(LEFT_MARGIN, Inches(6.2), CONTENT_WIDTH, Inches(0.5))
    tf = standards_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "STANDARDS: ICAO Doc 9303, ISO/IEC 14443, ISO 7816-4"
    p.font.name = BODY_FONT
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR

    cards_box = slide.shapes.add_textbox(LEFT_MARGIN, Inches(6.6), CONTENT_WIDTH, Inches(0.3))
    tf = cards_box.text_frame
    p = tf.paragraphs[0]
    p.text = "CARDS: Turkish eID, e-Passport, MIFARE, NDEF, 10+ types"
    p.font.name = BODY_FONT
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_GRAY

    add_page_number(slide, 10)


def create_slide_11_doc_demo(prs, layout):
    """SLIDE 11: Document Verification Demo"""
    slide = prs.slides.add_slide(layout)
    add_white_background(slide)

    add_title(slide, "DOCUMENT VERIFICATION IN ACTION")
    add_subtitle_line(slide, Inches(1.4))

    # Left section: Card Detection
    left_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        LEFT_MARGIN, Inches(1.8), Inches(5.8), Inches(4.2)
    )
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = LIGHT_GRAY
    left_box.line.fill.background()

    left_title = slide.shapes.add_textbox(LEFT_MARGIN + Inches(0.2), Inches(1.9), Inches(5.4), Inches(0.4))
    tf = left_title.text_frame
    p = tf.paragraphs[0]
    p.text = "CARD DETECTION"
    p.font.name = BODY_FONT
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR

    # Placeholder for card images
    card_placeholder = slide.shapes.add_textbox(LEFT_MARGIN + Inches(0.3), Inches(2.5), Inches(5.2), Inches(1.0))
    tf = card_placeholder.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "[TURKISH eID PHOTO]"
    p.font.name = BODY_FONT
    p.font.size = Pt(14)
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER

    # Detection results
    det_results = [
        ("Type: TURKISH_EID", SUCCESS_GREEN),
        ("Confidence: 97.3%", DARK_GRAY),
        ("CHECKMARK Ready for NFC", SUCCESS_GREEN)
    ]

    for i, (text, color) in enumerate(det_results):
        result_box = slide.shapes.add_textbox(LEFT_MARGIN + Inches(0.3), Inches(3.6 + i * 0.35), Inches(5.2), Inches(0.3))
        tf = result_box.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = BODY_FONT
        p.font.size = Pt(13)
        p.font.color.rgb = color

    passport_placeholder = slide.shapes.add_textbox(LEFT_MARGIN + Inches(0.3), Inches(4.8), Inches(5.2), Inches(0.8))
    tf = passport_placeholder.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "[E-PASSPORT PHOTO]"
    p.font.name = BODY_FONT
    p.font.size = Pt(14)
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = "Type: E_PASSPORT | Confidence: 98.1%"
    p.font.name = BODY_FONT
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER

    # Right section: NFC Reading
    right_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.9), Inches(1.8), Inches(5.8), Inches(4.2)
    )
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = LIGHT_GRAY
    right_box.line.fill.background()

    right_title = slide.shapes.add_textbox(Inches(7.1), Inches(1.9), Inches(5.4), Inches(0.4))
    tf = right_title.text_frame
    p = tf.paragraphs[0]
    p.text = "NFC CHIP READING"
    p.font.name = BODY_FONT
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR

    nfc_results = [
        "CHECKMARK BAC Authenticated",
        "CHECKMARK DG1 Personal Data",
        "CHECKMARK DG2 Photo Loaded",
        "CHECKMARK SOD Signature Valid"
    ]

    for i, result in enumerate(nfc_results):
        result_box = slide.shapes.add_textbox(Inches(7.3), Inches(2.5 + i * 0.45), Inches(5.2), Inches(0.4))
        tf = result_box.text_frame
        p = tf.paragraphs[0]
        p.text = result
        p.font.name = BODY_FONT
        p.font.size = Pt(14)
        p.font.color.rgb = SUCCESS_GREEN

    # DG2 photo placeholder
    dg2_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(7.3), Inches(4.5), Inches(5.0), Inches(1.3)
    )
    dg2_box.fill.solid()
    dg2_box.fill.fore_color.rgb = WHITE
    dg2_box.line.color.rgb = SUCCESS_GREEN
    dg2_box.line.width = Pt(2)

    dg2_text = slide.shapes.add_textbox(Inches(7.3), Inches(4.8), Inches(5.0), Inches(0.7))
    tf = dg2_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "[HIGH-RES PHOTO FROM DG2 CHIP]"
    p.font.name = BODY_FONT
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = "JPEG2000 Quality"
    p.font.name = BODY_FONT
    p.font.size = Pt(11)
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER

    # Caption
    caption = slide.shapes.add_textbox(LEFT_MARGIN, Inches(6.2), CONTENT_WIDTH, Inches(0.4))
    tf = caption.text_frame
    p = tf.paragraphs[0]
    p.text = "Left: ML model auto-detects card type | Right: NFC reads chip securely"
    p.font.name = BODY_FONT
    p.font.size = Pt(14)
    p.font.color.rgb = PRIMARY_COLOR
    p.alignment = PP_ALIGN.CENTER

    add_page_number(slide, 11)


def create_slide_12_tasks(prs, layout):
    """SLIDE 12: Tasks Accomplished"""
    slide = prs.slides.add_slide(layout)
    add_white_background(slide)

    add_title(slide, "WHAT WE BUILT - FALL 2025")
    add_subtitle_line(slide, Inches(1.4))

    # Column 1: Biometric Features
    col1_title = slide.shapes.add_textbox(LEFT_MARGIN, Inches(1.7), Inches(4.0), Inches(0.4))
    tf = col1_title.text_frame
    p = tf.paragraphs[0]
    p.text = "BIOMETRIC FEATURES"
    p.font.name = BODY_FONT
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR

    biometric_items = [
        "1. Face Detection (MediaPipe)",
        "2. Face Enrollment",
        "3. Face Verification (1:1)",
        "4. Face Search (1:N)",
        "5. Biometric Puzzle (Liveness)",
        "6. Frame Quality Analysis",
        "7. Demographic Analysis",
        "8. Similarity Scoring"
    ]

    for i, item in enumerate(biometric_items):
        item_box = slide.shapes.add_textbox(LEFT_MARGIN, Inches(2.1 + i * 0.35), Inches(4.0), Inches(0.3))
        tf = item_box.text_frame
        p = tf.paragraphs[0]
        p.text = item
        p.font.name = BODY_FONT
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_GRAY

    # Column 2: Document Verification
    col2_title = slide.shapes.add_textbox(Inches(4.5), Inches(1.7), Inches(4.0), Inches(0.4))
    tf = col2_title.text_frame
    p = tf.paragraphs[0]
    p.text = "DOCUMENT VERIFICATION"
    p.font.name = BODY_FONT
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR

    doc_items = [
        "9. Visual Card Detection (ML)",
        "10. NFC Document Reading",
        "11. BAC Authentication",
        "12. SOD Validation",
        "13. MRZ Parsing",
        "14. 10+ Card Types Support"
    ]

    for i, item in enumerate(doc_items):
        item_box = slide.shapes.add_textbox(Inches(4.5), Inches(2.1 + i * 0.35), Inches(4.0), Inches(0.3))
        tf = item_box.text_frame
        p = tf.paragraphs[0]
        p.text = item
        p.font.name = BODY_FONT
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_GRAY

    # Column 3: Infrastructure
    col3_title = slide.shapes.add_textbox(Inches(9.0), Inches(1.7), Inches(4.0), Inches(0.4))
    tf = col3_title.text_frame
    p = tf.paragraphs[0]
    p.text = "INFRASTRUCTURE"
    p.font.name = BODY_FONT
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR

    infra_items = [
        "15. JWT Authentication",
        "16. Multi-Tenant Architecture",
        "17. 40+ REST API Endpoints",
        "18. PostgreSQL + pgvector",
        "19. IVFFlat Vector Indexing",
        "20. Redis Caching",
        "21. Flyway Migrations",
        "22. Kotlin Multiplatform Apps"
    ]

    for i, item in enumerate(infra_items):
        item_box = slide.shapes.add_textbox(Inches(9.0), Inches(2.1 + i * 0.35), Inches(4.0), Inches(0.3))
        tf = item_box.text_frame
        p = tf.paragraphs[0]
        p.text = item
        p.font.name = BODY_FONT
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_GRAY

    # Web Demo GUI section
    gui_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        LEFT_MARGIN, Inches(5.3), Inches(4.0), Inches(1.5)
    )
    gui_box.fill.solid()
    gui_box.fill.fore_color.rgb = LIGHT_GRAY
    gui_box.line.fill.background()

    gui_title = slide.shapes.add_textbox(LEFT_MARGIN + Inches(0.2), Inches(5.4), Inches(3.6), Inches(0.3))
    tf = gui_title.text_frame
    p = tf.paragraphs[0]
    p.text = "WEB DEMO GUI"
    p.font.name = BODY_FONT
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR

    gui_items = ["Dashboard", "Enrollment Page", "Verification Page", "Search Page", "Liveness Testing", "Quality Analysis"]
    gui_text = slide.shapes.add_textbox(LEFT_MARGIN + Inches(0.2), Inches(5.7), Inches(3.6), Inches(1.0))
    tf = gui_text.text_frame
    tf.word_wrap = True
    for i, item in enumerate(gui_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.name = BODY_FONT
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_GRAY

    add_page_number(slide, 12)


def create_slide_13_challenges(prs, layout):
    """SLIDE 13: Technical Challenges & Solutions"""
    slide = prs.slides.add_slide(layout)
    add_white_background(slide)

    add_title(slide, "TECHNICAL CHALLENGES")
    add_subtitle_line(slide, Inches(1.4))

    challenges = [
        ("1. NFC PASSPORT INTEGRATION",
         "Complex ICAO protocols, BAC handshake, SOD parsing",
         "Modular reader architecture with 7 specialized readers"),
        ("2. CARD DETECTION MODEL TRAINING",
         "Dataset collection, model accuracy, real-time speed",
         "Custom dataset + optimized on-device inference"),
        ("3. CROSS-LANGUAGE MICROSERVICE COMMUNICATION",
         "Java <-> Python service integration, type safety",
         "REST APIs with OpenAPI contracts + Redis event bus")
    ]

    top = Inches(1.8)
    for title, challenge, solution in challenges:
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            LEFT_MARGIN, top, CONTENT_WIDTH, Inches(1.5)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GRAY
        box.line.fill.background()

        title_box = slide.shapes.add_textbox(LEFT_MARGIN + Inches(0.2), top + Inches(0.1), CONTENT_WIDTH, Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = BODY_FONT
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR

        challenge_box = slide.shapes.add_textbox(LEFT_MARGIN + Inches(0.2), top + Inches(0.5), CONTENT_WIDTH, Inches(0.4))
        tf = challenge_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"Challenge: {challenge}"
        p.font.name = BODY_FONT
        p.font.size = Pt(13)
        p.font.color.rgb = RgbColor(0xE7, 0x4C, 0x3C)  # Red

        solution_box = slide.shapes.add_textbox(LEFT_MARGIN + Inches(0.2), top + Inches(0.9), CONTENT_WIDTH, Inches(0.4))
        tf = solution_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"Solution: {solution}"
        p.font.name = BODY_FONT
        p.font.size = Pt(13)
        p.font.color.rgb = SUCCESS_GREEN

        top += Inches(1.7)

    add_page_number(slide, 13)


def create_slide_14_future(prs, layout):
    """SLIDE 14: Future Work & Contingency Plan"""
    slide = prs.slides.add_slide(layout)
    add_white_background(slide)

    add_title(slide, "SEMESTER 2 PLANS & B-PLAN")
    add_subtitle_line(slide, Inches(1.4))

    # Timeline section
    timeline_title = slide.shapes.add_textbox(LEFT_MARGIN, Inches(1.7), CONTENT_WIDTH, Inches(0.4))
    tf = timeline_title.text_frame
    p = tf.paragraphs[0]
    p.text = "SEMESTER 2 TIMELINE (Spring 2026)"
    p.font.name = BODY_FONT
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR

    timeline_items = [
        ("FEB", "Service Integration + OCR Module"),
        ("MAR", "Real-Time Proctoring Module"),
        ("APR", "Security Testing + Mobile Polish"),
        ("MAY", "Production Deployment")
    ]

    for i, (month, task) in enumerate(timeline_items):
        month_box = slide.shapes.add_textbox(LEFT_MARGIN + Inches(0.2), Inches(2.2 + i * 0.45), Inches(1.0), Inches(0.35))
        tf = month_box.text_frame
        p = tf.paragraphs[0]
        p.text = month
        p.font.name = BODY_FONT
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = ACCENT_COLOR

        arrow_box = slide.shapes.add_textbox(LEFT_MARGIN + Inches(1.2), Inches(2.2 + i * 0.45), Inches(0.5), Inches(0.35))
        tf = arrow_box.text_frame
        p = tf.paragraphs[0]
        p.text = ">>>"
        p.font.name = BODY_FONT
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY

        task_box = slide.shapes.add_textbox(LEFT_MARGIN + Inches(1.8), Inches(2.2 + i * 0.45), Inches(8.0), Inches(0.35))
        tf = task_box.text_frame
        p = tf.paragraphs[0]
        p.text = task
        p.font.name = BODY_FONT
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY

    # Contingency Plan box
    bplan_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        LEFT_MARGIN, Inches(4.3), CONTENT_WIDTH, Inches(1.1)
    )
    bplan_box.fill.solid()
    bplan_box.fill.fore_color.rgb = RgbColor(0xFF, 0xF3, 0xCD)  # Light yellow
    bplan_box.line.fill.background()

    bplan_title = slide.shapes.add_textbox(LEFT_MARGIN + Inches(0.2), Inches(4.4), CONTENT_WIDTH, Inches(0.3))
    tf = bplan_title.text_frame
    p = tf.paragraphs[0]
    p.text = "CONTINGENCY PLAN (B-PLAN)"
    p.font.name = BODY_FONT
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY

    bplan_items = [
        "IF NFC edge cases fail -> Focus on Turkish eID + e-Passport",
        "IF Integration delayed -> Web demo as primary deliverable"
    ]

    for i, item in enumerate(bplan_items):
        item_box = slide.shapes.add_textbox(LEFT_MARGIN + Inches(0.3), Inches(4.8 + i * 0.3), CONTENT_WIDTH, Inches(0.3))
        tf = item_box.text_frame
        p = tf.paragraphs[0]
        p.text = item
        p.font.name = BODY_FONT
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_GRAY

    # Future Research box
    future_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        LEFT_MARGIN, Inches(5.6), CONTENT_WIDTH, Inches(1.2)
    )
    future_box.fill.solid()
    future_box.fill.fore_color.rgb = LIGHT_GRAY
    future_box.line.fill.background()

    future_title = slide.shapes.add_textbox(LEFT_MARGIN + Inches(0.2), Inches(5.7), CONTENT_WIDTH, Inches(0.3))
    tf = future_title.text_frame
    p = tf.paragraphs[0]
    p.text = "FUTURE RESEARCH (Beyond Project)"
    p.font.name = BODY_FONT
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR

    future_items = [
        "Offline mode with on-device ML",
        "Fingerprint & iris biometrics",
        "Embedded devices (Raspberry Pi)"
    ]

    for i, item in enumerate(future_items):
        item_box = slide.shapes.add_textbox(LEFT_MARGIN + Inches(0.3), Inches(6.1 + i * 0.3), CONTENT_WIDTH, Inches(0.3))
        tf = item_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"• {item}"
        p.font.name = BODY_FONT
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_GRAY

    add_page_number(slide, 14)


def create_slide_15_thank_you(prs, layout):
    """SLIDE 15: Thank You"""
    slide = prs.slides.add_slide(layout)
    add_background(slide, PRIMARY_COLOR)

    # Thank you text
    thank_box = slide.shapes.add_textbox(LEFT_MARGIN, Inches(2.5), CONTENT_WIDTH, Inches(1.0))
    tf = thank_box.text_frame
    p = tf.paragraphs[0]
    p.text = "THANK YOU"
    p.font.name = TITLE_FONT
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Advisor acknowledgment
    advisor_box = slide.shapes.add_textbox(LEFT_MARGIN, Inches(3.8), CONTENT_WIDTH, Inches(1.0))
    tf = advisor_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "We thank our advisor, Assoc. Prof. Dr. Mustafa Agaoglu,"
    p.font.name = BODY_FONT
    p.font.size = Pt(20)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = "for his guidance throughout this project."
    p.font.name = BODY_FONT
    p.font.size = Pt(20)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER

    # GitHub link
    github_box = slide.shapes.add_textbox(LEFT_MARGIN, Inches(5.2), CONTENT_WIDTH, Inches(0.5))
    tf = github_box.text_frame
    p = tf.paragraphs[0]
    p.text = "GitHub: github.com/Rollingcat-Software/FIVUCSAS"
    p.font.name = BODY_FONT
    p.font.size = Pt(18)
    p.font.color.rgb = ACCENT_COLOR
    p.alignment = PP_ALIGN.CENTER

    # Open source note
    oss_box = slide.shapes.add_textbox(LEFT_MARGIN, Inches(5.7), CONTENT_WIDTH, Inches(0.4))
    tf = oss_box.text_frame
    p = tf.paragraphs[0]
    p.text = "(Will be released as open-source)"
    p.font.name = BODY_FONT
    p.font.size = Pt(14)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER

    # Page number (white for dark background)
    page_box = slide.shapes.add_textbox(Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.3))
    tf = page_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"15/{TOTAL_SLIDES}"
    p.font.name = BODY_FONT
    p.font.size = PAGE_NUM_SIZE
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.RIGHT


def create_slide_16_references(prs, layout):
    """SLIDE 16: References"""
    slide = prs.slides.add_slide(layout)
    add_white_background(slide)

    add_title(slide, "REFERENCES")
    add_subtitle_line(slide, Inches(1.4))

    references = [
        "[1] Taigman et al. (2014). DeepFace: Closing the Gap to Human-Level",
        "    Performance in Face Verification. CVPR.",
        "",
        "[2] Schroff et al. (2015). FaceNet: A Unified Embedding for Face",
        "    Recognition and Clustering. IEEE CVPR.",
        "",
        "[3] Deng et al. (2019). ArcFace: Additive Angular Margin Loss for",
        "    Deep Face Recognition. CVPR.",
        "",
        "[4] Lugaresi et al. (2019). MediaPipe: A Framework for Building",
        "    Perception Pipelines. Google Research.",
        "",
        "[5] ICAO Doc 9303: Machine Readable Travel Documents.",
        "",
        "[6] ISO/IEC 14443 & ISO 7816-4: NFC Standards."
    ]

    ref_box = slide.shapes.add_textbox(LEFT_MARGIN, Inches(1.8), CONTENT_WIDTH, Inches(5.0))
    tf = ref_box.text_frame
    tf.word_wrap = True

    for i, ref in enumerate(references):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = ref
        p.font.name = BODY_FONT
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(2)

    add_page_number(slide, 16)


def create_slide_17_qa(prs, layout):
    """SLIDE 17: Q&A"""
    slide = prs.slides.add_slide(layout)
    add_background(slide, PRIMARY_COLOR)

    # Q&A text
    qa_box = slide.shapes.add_textbox(LEFT_MARGIN, Inches(3.0), CONTENT_WIDTH, Inches(1.5))
    tf = qa_box.text_frame
    p = tf.paragraphs[0]
    p.text = "QUESTIONS & ANSWERS"
    p.font.name = TITLE_FONT
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(LEFT_MARGIN, Inches(4.5), CONTENT_WIDTH, Inches(0.5))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "We welcome your questions."
    p.font.name = BODY_FONT
    p.font.size = Pt(20)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER

    # Page number (white for dark background)
    page_box = slide.shapes.add_textbox(Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.3))
    tf = page_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"17/{TOTAL_SLIDES}"
    p.font.name = BODY_FONT
    p.font.size = PAGE_NUM_SIZE
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.RIGHT


# ============================================================================
# MAIN
# ============================================================================

def main():
    """Generate the presentation."""
    print("Generating FIVUCSAS Presentation...")
    print("=" * 50)

    prs = create_presentation()

    # Save the presentation
    output_path = os.path.join(os.path.dirname(__file__), "FIVUCSAS_Presentation.pptx")
    prs.save(output_path)

    print(f"Presentation saved to: {output_path}")
    print(f"Total slides: {TOTAL_SLIDES}")
    print("=" * 50)
    print("\nNOTE: Replace placeholder images with actual screenshots:")
    print("  - Slide 8: Team member faces with 468 mesh overlay")
    print("  - Slide 11: Card detection and NFC verification screenshots")
    print("  - Add Marmara University logo to Slide 1")
    print("\nDone!")


if __name__ == "__main__":
    main()
